from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = "Kubeflow_Workspace_Intelligence_LLMOps_Deck_v4.pptx"

INK = RGBColor(17, 24, 39)
MUTED = RGBColor(75, 85, 99)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
PINK = RGBColor(219, 39, 119)
TEAL = RGBColor(13, 148, 136)
PURPLE = RGBColor(124, 58, 237)
RED = RGBColor(220, 38, 38)
LIGHT = RGBColor(248, 250, 252)
PANEL = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)


deck = [
    {
        "kind": "title",
        "title": "Kubeflow Workspace Intelligence",
        "subtitle": "RAG-based chatbot with LLMOps observability and evaluation",
        "body": "Production-style assistant for workspace knowledge, notebook discovery, platform docs QA, and user expertise profiling.",
        "footer": "FastAPI | Next.js | Milvus | LiteLLM | Langfuse | RAGAS",
        "notes": "Frame this as a production-style LLMOps system, not a generic chatbot wrapper. The project connects ingestion, retrieval, generation, tracing, and evaluation into one inspectable system.",
    },
    {
        "kind": "bullets",
        "title": "Problem",
        "bullets": [
            "Kubeflow workspaces become knowledge silos as notebooks and scripts grow",
            "Filename search misses ML intent: tools, models, datasets, techniques",
            "Teams need grounded answers, expertise discovery, and operational visibility",
        ],
        "notes": "The pain is discovery plus trust. Notebooks hold institutional memory, but that memory is hard to query and hard to evaluate without tracing and quality metrics.",
    },
    {
        "kind": "bullets",
        "title": "Solution Overview",
        "bullets": [
            "Intent-routed RAG assistant over platform docs, artifact summaries, and user profiles",
            "Semantic retrieval backed by OpenAI embeddings and Milvus collections",
            "LiteLLM and Langfuse provide request-level visibility across classify, rewrite, and generate",
            "RAGAS and LLM-as-judge evaluate groundedness and relevance after the answer returns",
        ],
        "notes": "The important architecture choice is separating knowledge surfaces rather than putting everything into one mixed index. Each intent has a retrieval path that matches the question type.",
    },
    {
        "kind": "bullets",
        "title": "Dataset and Ingestion Scope",
        "bullets": [
            "Current checked-in catalog: 5 workspaces, 211 artifacts",
            "Artifact mix: 132 notebooks, 75 scripts, 4 text files",
            "CSV, binary, archive, and unsupported files are excluded from semantic indexing",
            "Full and incremental ingestion use content hashes to control reprocessing cost",
        ],
        "notes": "The original prompt mentions approximately 250 notebooks, but the actual catalog in this repo has 132 notebooks. In an interview, say the checked-in snapshot has 132 and the design scales to a larger corpus.",
    },
    {
        "kind": "bullets",
        "title": "Core Capabilities",
        "bullets": [
            "Workspace browsing and profiling",
            "Notebook/script semantic search",
            "LLM-generated artifact summaries and user profiles",
            "Platform documentation QA from Word documents",
            "Assistant with intent routing, query rewriting, sources, feedback, and trace IDs",
        ],
        "notes": "The implemented project goes beyond the original prompt. It includes platform-doc QA, hybrid routing, user resolution, feedback scoring, and admin sync endpoints.",
    },
    {
        "kind": "architecture",
        "title": "High-Level Architecture",
        "bullets": [
            "Next.js acts as a BFF and proxies to FastAPI",
            "FastAPI orchestrates search, chat, sync, metrics, and feedback",
            "Milvus stores separate knowledge collections",
            "LiteLLM and Langfuse make LLM calls observable",
        ],
        "notes": "Walk left to right: users access the Next.js UI; server-side routes call FastAPI; FastAPI orchestrates ingestion-backed retrieval, chat routing, LLM generation, tracing, and scoring.",
    },
    {
        "kind": "data_flow",
        "title": "Data Flow: Offline Build and Online Serving",
        "bullets": [
            "Offline path builds catalog, summaries, profiles, document chunks, and vector indexes",
            "Online path classifies the query, retrieves context, generates an answer, and returns trace_id",
            "Evaluation path posts quality scores to the same Langfuse trace asynchronously",
        ],
        "notes": "Use this slide to show that RAG has two lifecycles: building the knowledge layer and serving a user query. They meet at Milvus, and observability follows the online request.",
    },
    {
        "kind": "bullets",
        "title": "Ingestion and Indexing Pipeline",
        "bullets": [
            "Scan workspace directories and classify supported files",
            "Extract notebook/script metadata: tools, table references, database targets",
            "Convert artifacts into LangChain documents and extract notebook cell text",
            "Chunk by content type, embed with text-embedding-3-small, and upsert to Milvus",
        ],
        "notes": "Chunking is content-aware: recursive splitting for notebooks, code splitting for scripts, markdown splitting for markdown. Incremental indexing avoids re-embedding unchanged artifacts.",
    },
    {
        "kind": "bullets",
        "title": "Retrieval Surfaces",
        "bullets": [
            "kubeflow_artifacts: raw artifact chunks for semantic artifact search",
            "artifact_summaries: concise notebook/script descriptions and tags",
            "user_profiles: generated expertise summaries by workspace owner",
            "platform_docs: chunks extracted from platform Word documents",
        ],
        "notes": "Separate collections improve prompt clarity and retrieval precision. The tradeoff is more sync and freshness management.",
    },
    {
        "kind": "runtime",
        "title": "Chatbot Runtime Flow",
        "bullets": [
            "One trace_id groups classify, rewrite, retrieve, generate, and scoring",
            "Intent decides which store is queried and which prompt is built",
            "USER_SEARCH uses exact/candidate name resolution before semantic retrieval",
            "Scores and feedback are attached to the same Langfuse trace",
        ],
        "notes": "This is the core interview architecture slide. Emphasize that query classification is not decorative; it controls retrieval, cost, prompt shape, and failure handling.",
    },
    {
        "kind": "bullets",
        "title": "Intent Routing",
        "bullets": [
            "DOC_QA routes to platform_docs",
            "ARTIFACT_SEARCH routes to artifact_summaries",
            "USER_SEARCH routes to direct profile lookup or user_profiles",
            "HYBRID retrieves from docs, artifacts, and profiles",
            "OUT_OF_SCOPE returns a bounded response instead of general model knowledge",
        ],
        "notes": "The classifier returns intent and confidence. Failure falls back to DOC_QA with low confidence, which is conservative compared with generating ungrounded workspace facts.",
    },
    {
        "kind": "bullets",
        "title": "LLM Orchestration with LiteLLM",
        "bullets": [
            "All chat calls use an OpenAI-compatible client pointed at LiteLLM",
            "Generation defaults to gpt-4o-mini for cost-sensitive interaction",
            "Evaluation and judge paths use gpt-4o for stronger scoring reliability",
            "Trace metadata groups multi-step calls under one Langfuse request",
        ],
        "notes": "LiteLLM centralizes model routing, provider configuration, token/cost tracking, and Langfuse callbacks. Application code stays provider-agnostic.",
    },
    {
        "kind": "bullets",
        "title": "Layer 1 Observability",
        "bullets": [
            "LiteLLM forwards prompts, responses, latency, cost, tokens, and errors to Langfuse",
            "Application posts response_length, has_content, intent_confidence, and source_count",
            "The trace_id is returned so frontend feedback can attach to the same trace",
        ],
        "notes": "Layer 1 tells us what happened operationally: latency, cost, token use, prompt/response, errors, and lightweight response health signals.",
    },
    {
        "kind": "bullets",
        "title": "Layer 2 RAG Evaluation",
        "bullets": [
            "Runs after the response in a background thread",
            "RAGAS metrics: faithfulness, answer_relevancy, context_precision",
            "Exact USER_SEARCH uses LLM judge score: profile_relevance",
            "Scores are posted back to the same Langfuse trace",
        ],
        "notes": "Layer 2 tells us whether the answer was good and grounded. It is asynchronous so RAGAS does not add user-facing latency.",
    },
    {
        "kind": "bullets",
        "title": "Evaluation Interpretation",
        "bullets": [
            "Low faithfulness means answer claims are unsupported by retrieved context",
            "Low answer_relevancy means the answer missed the user question",
            "Low context_precision means retrieval returned weak or poorly ranked context",
            "Low profile_relevance means people-search output did not match the ask",
        ],
        "notes": "These metrics are diagnostic. They indicate whether to tune retrieval, chunking, prompts, classification, or profile generation.",
    },
    {
        "kind": "bullets",
        "title": "API and Web Surface",
        "bullets": [
            "/query, /chat, /workspaces, /user-profiles, /profile/workspace/{id}",
            "Admin endpoints rebuild indexes, summaries, profiles, and platform-doc chunks",
            "Feedback and score endpoints attach quality signals to Langfuse traces",
            "Next.js server routes proxy browser requests to FastAPI via PYTHON_API_URL",
        ],
        "notes": "The admin endpoints show lifecycle thinking. The system can rebuild and refresh its knowledge stores rather than being a one-off script.",
    },
    {
        "kind": "bullets",
        "title": "Design Decisions and Tradeoffs",
        "bullets": [
            "Separate vector collections improve precision but require freshness management",
            "Summary-first profiles reduce token noise but depend on summary quality",
            "Async evaluation preserves latency but scores arrive later",
            "LLM classifier is flexible but introduces uncertainty and model dependency",
            "LiteLLM adds a service but centralizes routing and observability",
        ],
        "notes": "A strong interview answer should not just list choices. It should describe the tradeoff and why the chosen side fits this system.",
    },
    {
        "kind": "bullets",
        "title": "Failure Cases and Mitigations",
        "bullets": [
            "Missing or stale indexes -> health checks and admin sync endpoints",
            "Low retrieval recall -> query rewriting, chunk overlap, and intent-specific stores",
            "Hallucinated answers -> OUT_OF_SCOPE guard, grounded prompts, faithfulness scoring",
            "Ambiguous people queries -> candidate resolution before semantic fallback",
            "Evaluation cost/latency -> background RAGAS execution",
        ],
        "notes": "Name realistic RAG failures and tie each one to a detection or mitigation mechanism. This makes the architecture feel production-minded.",
    },
    {
        "kind": "bullets",
        "title": "Production Readiness: Current State",
        "bullets": [
            "Implemented: full/incremental ingestion, indexing, assistant runtime, tracing, scoring, RAGAS, feedback, admin sync",
            "Observable request lifecycle: classify, rewrite, retrieve, generate, score",
            "Operational controls exist, but identity, ACLs, alerting, and release gates still need hardening",
        ],
        "notes": "Be honest: this is production-style now. It becomes production-ready after the security, reliability, deployment, and quality gates are added.",
    },
    {
        "kind": "bullets",
        "title": "Production Readiness: Security and Governance",
        "bullets": [
            "Add authentication and propagate user identity from UI to FastAPI",
            "Apply workspace/user ACL filters before retrieval and before answer generation",
            "Move secrets to a managed secret store and rotate LiteLLM/Langfuse keys",
            "Add audit logs for admin syncs, data access, feedback, and prompt changes",
            "Introduce PII/sensitive-content checks for artifacts and trace payloads",
        ],
        "notes": "This is the first production-readiness path: prevent data leakage. In enterprise RAG, ACL-aware retrieval is not optional because generated answers can expose retrieved context.",
    },
    {
        "kind": "bullets",
        "title": "Production Readiness: Reliability and Quality Gates",
        "bullets": [
            "Move background RAGAS work to a queue-backed worker for retry and backpressure",
            "Create scheduled eval datasets for DOC_QA, ARTIFACT_SEARCH, USER_SEARCH, and HYBRID",
            "Gate prompt/model/index changes on regression metrics before promotion",
            "Add SLO dashboards for latency, error rate, index freshness, cost/query, and faithfulness",
            "Use canary releases for classifier prompts and generation prompts",
        ],
        "notes": "This is how production readiness can be met operationally: make quality measurable, make evaluation reliable, and make changes releasable with regression checks.",
    },
    {
        "kind": "bullets",
        "title": "Production Readiness: Deployment and Operations",
        "bullets": [
            "Containerize FastAPI, Next.js, LiteLLM, Langfuse, and Milvus with environment-specific configs",
            "Run ingestion/indexing as scheduled jobs with idempotent full and incremental modes",
            "Add CI for unit, integration, retrieval, and API contract tests",
            "Define runbooks for stale indexes, failed LLM calls, high cost, and low evaluation scores",
        ],
        "notes": "This slide completes the production-readiness path. It explains how the system moves from local/demo operation to repeatable environments and on-call ownership.",
    },
    {
        "kind": "bullets",
        "title": "Demo Walkthrough",
        "bullets": [
            "Show workspaces and a generated workspace/user profile",
            "Run semantic artifact search for an ML concept",
            "Ask DOC_QA, ARTIFACT_SEARCH, and USER_SEARCH questions in the assistant",
            "Open Langfuse trace and show classify, rewrite, generate, and score events",
        ],
        "notes": "Demo the loop: user question, routed answer, trace, quality score, diagnosis. Keep the demo focused on engineering behavior rather than UI decoration.",
    },
    {
        "kind": "bullets",
        "title": "Close",
        "bullets": [
            "Turns workspace artifacts into a queryable intelligence layer",
            "Separates ingestion, retrieval, orchestration, observability, and evaluation",
            "Built to explain failures, not just return answers",
            "Clear production path: security, quality gates, deployment automation, and operational runbooks",
        ],
        "notes": "The strongest takeaway is engineering discipline: observable, evaluable, grounded RAG with a practical hardening path.",
    },
]


def add_textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def set_text(shape, text, size=12, color=INK, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.07)
    tf.margin_right = Inches(0.07)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold


def add_box(slide, x, y, w, h, text, fill, line=WHITE, size=12, bold=False, color=INK):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    set_text(shape, text, size=size, bold=bold, color=color)
    return shape


def add_label(slide, x, y, w, h, text, fill=PANEL, color=INK):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = WHITE
    set_text(shape, text, size=10, bold=True, color=color)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_bullets_panel(slide, bullets, x=10.25, y=1.35, w=2.55, h=4.95):
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = RGBColor(226, 232, 240)
    add_textbox(slide, x + 0.18, y + 0.2, w - 0.35, 0.28, "Important Points", 12, INK, True)
    yy = y + 0.68
    for bullet in bullets:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22), Inches(yy + 0.08), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        add_textbox(slide, x + 0.4, yy, w - 0.55, 0.5, bullet, 10.5, INK)
        yy += 0.78


def add_brand(slide, index):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    for i, color in enumerate([BLUE, GREEN, AMBER, PINK]):
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i * 3.33), Inches(0.32), Inches(3.33), Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
    add_textbox(
        slide,
        0.65,
        7.08,
        12.0,
        0.28,
        f"Kubeflow Workspace Intelligence LLMOps - v4 - {index:02d}",
        size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def add_title(slide, title):
    add_textbox(slide, 0.65, 0.72, 12.0, 0.55, title, 30, INK, True)


def make_title(prs, item, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_brand(slide, index)
    add_textbox(slide, 0.8, 1.45, 11.8, 0.85, item["title"], 38, INK, True)
    add_textbox(slide, 0.82, 2.38, 11.2, 0.55, item["subtitle"], 22, BLUE)
    add_textbox(slide, 0.82, 3.25, 10.9, 1.1, item["body"], 20, MUTED)
    add_textbox(slide, 0.82, 5.45, 11.5, 0.4, item["footer"], 17, INK, True)
    add_notes(slide, item["notes"])


def make_bullet(prs, item, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_brand(slide, index)
    add_title(slide, item["title"])

    y = 1.52
    max_bullets = len(item["bullets"])
    step = 0.68 if max_bullets <= 4 else 0.58
    size = 19 if max_bullets <= 4 else 17
    for bullet in item["bullets"]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.86), Inches(y + 0.11), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        add_textbox(slide, 1.12, y, 11.35, 0.55, bullet, size, INK)
        y += step
    add_notes(slide, item["notes"])


def make_architecture(prs, item, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_brand(slide, index)
    add_title(slide, item["title"])

    add_label(slide, 0.75, 1.35, 8.95, 0.28, "Request Path")
    add_box(slide, 0.75, 1.85, 1.35, 0.72, "Users\nDS / MLE", WHITE, RGBColor(203, 213, 225), 11, True)
    add_box(slide, 2.45, 1.85, 1.45, 0.72, "Next.js\nBFF", RGBColor(219, 234, 254), BLUE, 11, True)
    add_box(slide, 4.25, 1.85, 1.55, 0.72, "FastAPI\nRetrieval API", RGBColor(220, 252, 231), GREEN, 10.5, True)
    add_box(slide, 6.15, 1.85, 1.75, 0.72, "ChatEngine\norchestrator", RGBColor(254, 243, 199), AMBER, 10.5, True)
    add_box(slide, 8.25, 1.85, 1.45, 0.72, "Response\n+ trace_id", WHITE, RGBColor(203, 213, 225), 10.5, True)
    for x1, x2 in [(2.1, 2.45), (3.9, 4.25), (5.8, 6.15), (7.9, 8.25)]:
        add_arrow(slide, x1, 2.21, x2, 2.21, BLUE)

    add_label(slide, 0.75, 3.15, 4.3, 0.28, "Knowledge Layer")
    add_box(slide, 0.75, 3.6, 1.85, 0.62, "kubeflow_artifacts\nraw chunks", RGBColor(255, 247, 237), AMBER, 9.5)
    add_box(slide, 2.8, 3.6, 1.85, 0.62, "artifact_summaries\nnotebook summaries", RGBColor(255, 247, 237), AMBER, 9.5)
    add_box(slide, 0.75, 4.45, 1.85, 0.62, "user_profiles\nexpertise vectors", RGBColor(255, 247, 237), AMBER, 9.5)
    add_box(slide, 2.8, 4.45, 1.85, 0.62, "platform_docs\ndoc chunks", RGBColor(255, 247, 237), AMBER, 9.5)
    add_box(slide, 1.85, 5.32, 1.75, 0.42, "Milvus", RGBColor(254, 243, 199), AMBER, 12, True)

    add_label(slide, 5.4, 3.15, 4.3, 0.28, "LLMOps Layer")
    add_box(slide, 5.4, 3.6, 1.65, 0.68, "LiteLLM\nproxy", RGBColor(237, 233, 254), PURPLE, 11, True)
    add_box(slide, 7.3, 3.6, 1.65, 0.68, "Model\nprovider", RGBColor(237, 233, 254), PURPLE, 11, True)
    add_box(slide, 5.4, 4.55, 1.65, 0.68, "Langfuse\ntraces + scores", RGBColor(252, 231, 243), PINK, 10.5, True)
    add_box(slide, 7.3, 4.55, 1.65, 0.68, "RAGAS\nLLM judge", RGBColor(252, 231, 243), PINK, 10.5, True)
    add_arrow(slide, 7.0, 3.94, 7.3, 3.94, PURPLE)
    add_arrow(slide, 6.22, 4.28, 6.22, 4.55, PINK)
    add_arrow(slide, 7.3, 4.89, 7.05, 4.89, PINK)
    add_arrow(slide, 6.9, 2.57, 2.75, 3.6, AMBER)
    add_arrow(slide, 6.95, 2.57, 6.2, 3.6, PURPLE)

    add_bullets_panel(slide, item["bullets"])
    add_notes(slide, item["notes"])


def make_data_flow(prs, item, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_brand(slide, index)
    add_title(slide, item["title"])

    add_label(slide, 0.65, 1.35, 9.35, 0.28, "Offline Build Path")
    offline = [
        (0.65, "Workspace\nfiles"),
        (2.12, "Ingestion\ncatalog"),
        (3.6, "Parse +\nchunk"),
        (5.08, "Embeddings"),
        (6.55, "Milvus\ncollections"),
        (8.03, "Summaries\nprofiles docs"),
    ]
    colors = [WHITE, RGBColor(220, 252, 231), RGBColor(219, 234, 254), RGBColor(237, 233, 254), RGBColor(254, 243, 199), RGBColor(255, 247, 237)]
    for idx, (x, label) in enumerate(offline):
        add_box(slide, x, 1.88, 1.18, 0.72, label, colors[idx], RGBColor(203, 213, 225), 9.5, True)
        if idx < len(offline) - 1:
            add_arrow(slide, x + 1.18, 2.24, offline[idx + 1][0], 2.24, GREEN)

    add_label(slide, 0.65, 3.25, 9.35, 0.28, "Online Query Path")
    online = [
        (0.65, "User\nquery"),
        (2.12, "Classify\nrewrite"),
        (3.6, "Retrieve\ncontext"),
        (5.08, "Prompt +\ngenerate"),
        (6.55, "Answer +\nsources"),
        (8.03, "trace_id\nfeedback"),
    ]
    online_colors = [WHITE, RGBColor(220, 252, 231), RGBColor(254, 243, 199), RGBColor(237, 233, 254), RGBColor(219, 234, 254), RGBColor(252, 231, 243)]
    for idx, (x, label) in enumerate(online):
        add_box(slide, x, 3.78, 1.18, 0.72, label, online_colors[idx], RGBColor(203, 213, 225), 9.5, True)
        if idx < len(online) - 1:
            add_arrow(slide, x + 1.18, 4.14, online[idx + 1][0], 4.14, BLUE)

    add_arrow(slide, 7.14, 2.6, 4.15, 3.78, AMBER)
    add_textbox(slide, 5.05, 2.75, 2.0, 0.32, "retrieval uses built indexes", 9, MUTED, False, PP_ALIGN.CENTER)

    add_label(slide, 0.65, 5.15, 9.35, 0.28, "Evaluation and Improvement Path")
    add_box(slide, 1.0, 5.62, 1.55, 0.58, "Response\ncontexts", RGBColor(219, 234, 254), BLUE, 9.5, True)
    add_box(slide, 3.0, 5.62, 1.55, 0.58, "RAGAS /\nLLM judge", RGBColor(252, 231, 243), PINK, 9.5, True)
    add_box(slide, 5.0, 5.62, 1.55, 0.58, "Langfuse\nscores", RGBColor(252, 231, 243), PINK, 9.5, True)
    add_box(slide, 7.0, 5.62, 1.55, 0.58, "Tune prompts\nretrieval", RGBColor(220, 252, 231), GREEN, 9.5, True)
    for x1, x2 in [(2.55, 3.0), (4.55, 5.0), (6.55, 7.0)]:
        add_arrow(slide, x1, 5.91, x2, 5.91, PINK)

    add_bullets_panel(slide, item["bullets"])
    add_notes(slide, item["notes"])


def make_runtime(prs, item, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_brand(slide, index)
    add_title(slide, item["title"])

    main = [
        (0.65, 1.72, "Query\n+ trace_id", WHITE, BLUE),
        (2.0, 1.72, "Intent\nclassifier", RGBColor(220, 252, 231), GREEN),
        (3.35, 1.72, "Query\nrewriter", RGBColor(220, 252, 231), GREEN),
        (4.7, 1.72, "Routed\nretriever", RGBColor(254, 243, 199), AMBER),
        (6.05, 1.72, "Prompt\nbuilder", RGBColor(219, 234, 254), BLUE),
        (7.4, 1.72, "LiteLLM\ngenerate", RGBColor(237, 233, 254), PURPLE),
        (8.75, 1.72, "Formatted\nresponse", WHITE, BLUE),
    ]
    for idx, (x, y, label, fill, line) in enumerate(main):
        add_box(slide, x, y, 1.05, 0.72, label, fill, line, 9.2, True)
        if idx < len(main) - 1:
            add_arrow(slide, x + 1.05, y + 0.36, main[idx + 1][0], y + 0.36, BLUE)

    add_label(slide, 0.65, 3.05, 4.05, 0.28, "USER_SEARCH special path")
    add_box(slide, 0.95, 3.55, 1.25, 0.62, "Candidate\nresolution", RGBColor(220, 252, 231), GREEN, 9.5, True)
    add_box(slide, 2.55, 3.55, 1.25, 0.62, "Exact match:\nprofile fetch", RGBColor(254, 243, 199), AMBER, 9.2, True)
    add_arrow(slide, 2.2, 3.86, 2.55, 3.86, GREEN)
    add_arrow(slide, 2.53, 2.44, 1.57, 3.55, GREEN)
    add_arrow(slide, 3.18, 4.17, 8.95, 2.44, AMBER)

    add_label(slide, 5.05, 3.05, 4.75, 0.28, "Observability and evaluation")
    add_box(slide, 5.25, 3.55, 1.25, 0.62, "Layer 1\nheuristics", RGBColor(252, 231, 243), PINK, 9.5, True)
    add_box(slide, 6.85, 3.55, 1.25, 0.62, "Layer 2\nRAGAS", RGBColor(252, 231, 243), PINK, 9.5, True)
    add_box(slide, 8.45, 3.55, 1.25, 0.62, "Langfuse\ntrace", RGBColor(252, 231, 243), PINK, 9.5, True)
    add_arrow(slide, 6.5, 3.86, 6.85, 3.86, PINK)
    add_arrow(slide, 8.1, 3.86, 8.45, 3.86, PINK)
    add_arrow(slide, 9.28, 2.44, 5.88, 3.55, PINK)

    add_label(slide, 0.65, 5.08, 9.15, 0.28, "Intent-to-store routing")
    add_box(slide, 0.95, 5.55, 1.5, 0.52, "DOC_QA\nplatform_docs", RGBColor(255, 247, 237), AMBER, 8.8)
    add_box(slide, 2.75, 5.55, 1.5, 0.52, "ARTIFACT_SEARCH\nsummaries", RGBColor(255, 247, 237), AMBER, 8.5)
    add_box(slide, 4.55, 5.55, 1.5, 0.52, "USER_SEARCH\nprofiles", RGBColor(255, 247, 237), AMBER, 8.8)
    add_box(slide, 6.35, 5.55, 1.5, 0.52, "HYBRID\nall stores", RGBColor(255, 247, 237), AMBER, 8.8)
    add_box(slide, 8.15, 5.55, 1.5, 0.52, "OUT_OF_SCOPE\nbounded reply", RGBColor(254, 226, 226), RED, 8.8)

    add_bullets_panel(slide, item["bullets"])
    add_notes(slide, item["notes"])


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for idx, item in enumerate(deck, start=1):
        if item["kind"] == "title":
            make_title(prs, item, idx)
        elif item["kind"] == "architecture":
            make_architecture(prs, item, idx)
        elif item["kind"] == "data_flow":
            make_data_flow(prs, item, idx)
        elif item["kind"] == "runtime":
            make_runtime(prs, item, idx)
        else:
            make_bullet(prs, item, idx)

    prs.save(OUTPUT)
    print(f"Created {OUTPUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
