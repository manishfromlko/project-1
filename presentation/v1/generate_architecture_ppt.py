from pptx import Presentation
from pptx.util import Inches, Pt


TITLE = "Kubeflow Workspace Profiling Platform"
SUBTITLE = "Technical Architecture Deep Dive for Senior Management"


slides_content = [
    ("Executive Context", [
        "Platform objective: enterprise-grade GenAI assistant for workspace intelligence.",
        "Converges retrieval, profiling, and conversational analytics into one stack.",
        "Designed for scale: modular pipelines, vector-native storage, prompt externalization.",
    ]),
    ("Business Problem Framing", [
        "Teams struggle to discover reusable assets, experts, and operational know-how.",
        "Knowledge exists across notebooks/scripts/docs but remains fragmented.",
        "Solution creates a continuously indexed and queryable knowledge fabric.",
    ]),
    ("Platform Value Proposition", [
        "Semantic artifact discovery across notebooks, scripts, and markdown/text.",
        "LLM-generated artifact summaries and user expertise profiles.",
        "Intent-aware enterprise chatbot with source-grounded responses.",
    ]),
    ("System Overview", [
        "Frontend: Next.js webapp with App Router and React Query.",
        "Backend: FastAPI orchestration layer for search, profiling, chat, and admin sync.",
        "Data/AI plane: Milvus vector collections + OpenAI embeddings/LLM calls.",
    ]),
    ("Reference Architecture (Layered)", [
        "Experience layer: dashboard, search, workspace details, chat panel.",
        "Application layer: API proxy routes + FastAPI endpoint orchestration.",
        "Intelligence layer: retrieval services, classifiers, rewriters, prompt pipelines.",
        "Data layer: ingestion catalog, Milvus collections, platform docs corpus.",
    ]),
    ("Ingestion Pipeline Architecture", [
        "IngestionPipeline scans workspace folders and classifies artifact types.",
        "Guardrails block sensitive/unsupported files before indexing.",
        "Outputs normalized catalog and audit trail for downstream indexers.",
    ]),
    ("Incremental Processing Mechanics", [
        "Content hashes detect NEW / UPDATED / UNCHANGED artifact states.",
        "Incremental mode indexes only deltas to reduce embedding spend.",
        "Airflow DAG schedules daily incremental refresh for operational continuity.",
    ]),
    ("Vector Retrieval Architecture", [
        "Indexer converts catalog artifacts into embeddings and inserts into Milvus.",
        "Hybrid and vector retrievers support semantic discovery across artifacts.",
        "Response contracts include score + metadata for downstream UX rendering.",
    ]),
    ("Generative Enrichment Layer", [
        "Artifact Summary Indexer creates compact LLM summaries + tags.",
        "Profile Indexer aggregates summaries to generate per-user expertise vectors.",
        "Enrichment layer shifts from raw-code recall to semantic explainability.",
    ]),
    ("Chat Engine Orchestration", [
        "ChatEngine pipeline: classify -> rewrite -> retrieve -> generate -> format.",
        "Intent classes: DOC_QA, ARTIFACT_SEARCH, USER_SEARCH, HYBRID.",
        "USER_SEARCH uses deterministic name resolution before vector fallback.",
    ]),
    ("PromptOps and Governance", [
        "Prompts fully externalized under prompts/; no inline hardcoded prompt text.",
        "Enables fast policy/behavior updates without app code deployment.",
        "Supports model experimentation while preserving orchestration logic.",
    ]),
    ("Frontend Architecture", [
        "AppShell governs layout with collapsible AI assistant panel.",
        "React Query standardizes caching/retries/staleness for API interactions.",
        "Next.js API routes behave as BFF adapters and response normalizers.",
    ]),
    ("Backend API Surface", [
        "Core APIs: /query, /workspaces, /profile/workspace/{id}, /chat.",
        "Admin APIs trigger sync pipelines for artifacts, profiles, docs ingestion.",
        "Health/metrics endpoints expose runtime reliability and utilization indicators.",
    ]),
    ("Data Model and Vector Collections", [
        "kubeflow_artifacts: chunk vectors for semantic artifact retrieval.",
        "artifact_summaries: LLM summaries + tags for explainable discovery.",
        "user_profiles: expertise embeddings for people finder.",
        "platform_docs: doc chunks for enterprise documentation QA.",
    ]),
    ("End-to-End Data Flow", [
        "Raw files -> ingestion catalog -> vector indexing -> user query response.",
        "Chat path dynamically routes retrieval sources based on intent.",
        "Output includes grounded artifacts/users/sources for trust and traceability.",
    ]),
    ("Runtime Sequence: Search", [
        "UI emits SearchQuery -> Next API /api/search -> FastAPI /query.",
        "Retriever executes vector/hybrid lookup against kubeflow_artifacts.",
        "Result normalization in BFF maps backend payload to frontend cards.",
    ]),
    ("Runtime Sequence: Chat", [
        "UI sends query/history -> /api/chat -> FastAPI /chat endpoint.",
        "ChatEngine executes intent classification and query rewriting.",
        "Retriever fan-out and constrained generation return structured answer.",
    ]),
    ("Reliability and Observability", [
        "Startup preloads collections and degrades gracefully if optional stores missing.",
        "Metrics track uptime, query count, latency, memory, and error rate.",
        "Health endpoint reports vector connectivity + model readiness + cache stats.",
    ]),
    ("Security and Risk Controls", [
        "Guardrails skip sensitive/unsupported artifacts before downstream processing.",
        "Server-side API proxy keeps backend URL hidden from browsers.",
        "Risks: external LLM dependency, prompt drift, vector quality drift over time.",
    ]),
    ("Scalability and Performance Posture", [
        "Incremental indexing controls compute cost and ingestion latency.",
        "Collection separation isolates workloads by domain and intent.",
        "Bottlenecks to monitor: embedding throughput, Milvus query latency, token cost.",
    ]),
    ("Roadmap Alignment", [
        "Completed: ingestion, retrieval API, frontend, summaries/profiles, chatbot.",
        "Planned: LiteLLM gateway, Langfuse observability, production deployment hardening.",
        "Near-term focus: policy controls, tracing, eval harness, SLO-backed operations.",
    ]),
    ("Management Takeaways", [
        "Architecture is production-leaning with clear modular boundaries.",
        "Strong foundation for enterprise GenAI search + expert discovery use cases.",
        "Recommended next step: harden MLOps/LLMOps guardrails and observability.",
    ]),
]


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)


def main() -> None:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = TITLE
    title_slide.placeholders[1].text = SUBTITLE

    for title, bullets in slides_content:
        add_bullets_slide(prs, title, bullets)

    # Add a final appendix slide to reach 24 total pages
    appendix = prs.slides.add_slide(prs.slide_layouts[5])
    appendix.shapes.title.text = "Appendix: Key Module Map"
    tx_box = appendix.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    module_lines = [
        "src/ingestion/: pipeline.py, guards.py, extractors.py, storage.py, cli.py",
        "src/retrieval/: api.py, indexer.py, retriever.py, vector_store.py, embeddings.py",
        "src/retrieval/chatbot/: engine.py, classifier.py, query_rewriter.py, retrievers.py",
        "webapp/app/api/: BFF proxy routes for backend endpoint mediation",
        "airflow/dags/ingestion_dag.py: scheduled incremental ingestion orchestration",
    ]
    for idx, line in enumerate(module_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)

    prs.save("Kubeflow_Workspace_Profiling_Architecture_Deck.pptx")
    print("Created Kubeflow_Workspace_Profiling_Architecture_Deck.pptx with", len(prs.slides), "slides")


if __name__ == "__main__":
    main()
