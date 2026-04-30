from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = "Databricks_Intelligence_Layer_Azure_Architecture_v3.pptx"

NAVY = RGBColor(11, 31, 95)
BLUE = RGBColor(31, 78, 170)
CYAN = RGBColor(0, 163, 224)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(239, 245, 255)
TEXT = RGBColor(24, 31, 52)


slides = [
    ("Executive Summary", [
        "Objective: replicate Kubeflow intelligence layer in Azure Databricks with enterprise-grade operations.",
        "Core capability set remains: artifact discovery, user profiling, docs QA, and intent-routed assistant.",
        "Recommendation: hybrid architecture with Databricks-native data/AI plane plus Azure-hosted app/API plane.",
    ]),
    ("Current-to-Target Mapping", [
        "Filesystem ingestion -> Databricks Workspace/Repos/Jobs metadata ingestion.",
        "Milvus vector store -> Mosaic AI Vector Search indexes.",
        "Airflow orchestration -> Databricks Workflows with job dependencies and schedules.",
    ]),
    ("Target Platform Capabilities", [
        "Semantic artifact retrieval over user workspaces and repositories.",
        "Generated artifact summaries and recency-aware user expertise profiles.",
        "Enterprise assistant with citations and ACL-aware response filtering.",
    ]),
    ("Reference Architecture", [
        "Experience layer: Next.js UI with secure BFF API routes.",
        "Orchestration layer: API service for intent routing and response composition.",
        "Intelligence layer: Databricks serving endpoints + vector retrieval.",
        "Data layer: Unity Catalog governed Delta tables and indexes.",
    ]),
    ("Databricks Data Architecture", [
        "Bronze: raw metadata snapshots from workspace/repo/job APIs.",
        "Silver: normalized artifact registry with ACL and lineage metadata.",
        "Gold: summaries, profiles, and serving marts for assistant workloads.",
    ]),
    ("Ingestion and Cataloging Design", [
        "Metadata crawler runs as Databricks Workflow task chain.",
        "Change detection based on hash/version to support differential processing.",
        "Idempotent writes to Delta ensure replay-safe operations.",
    ]),
    ("Vectorization and Index Strategy", [
        "Chunking strategy tuned per artifact type (notebook/script/doc).",
        "Embeddings generated via Databricks-supported model endpoints or Azure OpenAI.",
        "Separate indexes: artifact chunks, summaries, profiles, documentation.",
    ]),
    ("Assistant Runtime Flow", [
        "Intent classification routes to DOC_QA, ARTIFACT_SEARCH, USER_SEARCH, HYBRID.",
        "Query rewrite improves recall and normalizes user intent vocabulary.",
        "Context assembly + constrained generation return grounded and cited responses.",
    ]),
    ("Access Control Model", [
        "User identity federated through Microsoft Entra ID.",
        "Retrieval layer applies workspace/team ACL filters before generation.",
        "Unity Catalog enforces table/index level governance and lineage traceability.",
    ]),
    ("Deployment Topology on Azure", [
        "Edge: Azure Front Door + WAF for ingress and policy enforcement.",
        "App plane: Next.js + Orchestrator API on App Service or AKS.",
        "Data/AI plane: Azure Databricks with private networking and endpoints.",
    ]),
    ("Network and Security Controls", [
        "Private Link to Databricks, Azure OpenAI, and ADLS dependencies.",
        "Secrets centralized in Azure Key Vault with managed identities.",
        "Audit logs and SIEM integration for compliance monitoring.",
    ]),
    ("CI/CD and Environment Strategy", [
        "Three-environment model: dev, staging, prod with isolated workspaces/catalogs.",
        "Databricks Asset Bundles for jobs/pipelines/config deployment automation.",
        "GitHub Actions/Azure DevOps for coordinated API, UI, and data-pipeline releases.",
    ]),
    ("Observability Blueprint", [
        "Platform telemetry: job runtime, freshness SLA, index health, endpoint latency.",
        "Application telemetry: request traces, intent mix, failure taxonomies.",
        "Business telemetry: answer quality, citation coverage, user adoption, cost per query.",
    ]),
    ("Reliability Engineering", [
        "SLOs for assistant latency, index freshness, and API availability.",
        "Dead-letter workflows for parsing/enrichment failures with replay support.",
        "Graceful degradation to retrieval-only responses when generation dependencies fail.",
    ]),
    ("FinOps Controls", [
        "Incremental indexing to reduce embedding and compute spend.",
        "Embedding cache and context budget constraints per query type.",
        "Unit economics dashboard for token cost, compute cost, and feature ROI.",
    ]),
    ("Data Quality and Evaluation", [
        "Offline eval sets for retrieval precision@k and groundedness.",
        "Prompt/version tracking with regression checks before promotion.",
        "Human feedback loop to tune ranking, summarization, and intent logic.",
    ]),
    ("Operating Model and Ownership", [
        "Data platform team owns pipelines, governance, and SLA reporting.",
        "AI platform team owns model endpoints, prompts, and evaluation lifecycle.",
        "Product engineering team owns UI/API experience and feature roadmap.",
    ]),
    ("Migration Plan (90 Days)", [
        "Phase 1: foundational ingestion, vector indexes, baseline search endpoints.",
        "Phase 2: summaries/profiles generation and assistant orchestration integration.",
        "Phase 3: hardening, observability, security controls, and production go-live gates.",
    ]),
    ("Risks and Mitigations", [
        "Risk: schema drift in source metadata -> mitigation: contract tests + schema registry.",
        "Risk: LLM variance/hallucination -> mitigation: retrieval grounding and strict formatting.",
        "Risk: cost spikes -> mitigation: quotas, caching, and differential indexing.",
    ]),
    ("Decision Ask for Leadership", [
        "Approve hybrid deployment model with Databricks-native intelligence core.",
        "Fund observability + governance as first-class capabilities, not post-go-live add-ons.",
        "Authorize phased rollout with measurable quality, reliability, and cost KPIs.",
    ]),
]


def add_brand(slide):
    top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.35))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0.35), Inches(13.33), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35))
    tf = footer.text_frame
    tf.text = "Azure Databricks Intelligence Layer - Ops Architecture"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = BLUE


def make_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_brand(slide)
    slide.shapes.title.text = "Databricks Intelligence Layer"
    t = slide.shapes.title.text_frame.paragraphs[0]
    t.font.size = Pt(42)
    t.font.bold = True
    t.font.color.rgb = NAVY
    slide.placeholders[1].text = "Azure Deployment and Operations Blueprint"
    s = slide.placeholders[1].text_frame.paragraphs[0]
    s.font.size = Pt(23)
    s.font.color.rgb = BLUE


def make_bullet(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_brand(slide)
    slide.shapes.title.text = title
    tp = slide.shapes.title.text_frame.paragraphs[0]
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = NAVY

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(21)
        p.font.color.rgb = TEXT
        p.level = 0


def main():
    prs = Presentation()
    make_title(prs)
    for title, bullets in slides:
        make_bullet(prs, title, bullets)
    prs.save(OUTPUT)
    print(f"Created {OUTPUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
