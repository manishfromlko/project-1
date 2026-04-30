from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE = "Kubeflow Workspace Profiling Platform"
SUBTITLE = "Reliance Technical Architecture Briefing (GenAI Program)"
OUTPUT_FILE = "Kubeflow_Workspace_Profiling_Architecture_Deck_v2_Reliance.pptx"

# Reliance-inspired corporate palette (approximation)
NAVY = RGBColor(11, 31, 95)
BLUE = RGBColor(31, 78, 170)
CYAN = RGBColor(0, 163, 224)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(25, 33, 56)
LIGHT_BG = RGBColor(238, 244, 255)


slides_content = [
    ("Executive Narrative", [
        "Strategic objective: convert distributed DS artifacts into a governed GenAI knowledge graph.",
        "Platform unifies semantic retrieval, profile intelligence, and conversational reasoning.",
        "Design favors modularity, incremental processing, and enterprise observability readiness.",
    ]),
    ("Why This Matters for Reliance", [
        "Accelerates AI asset reuse across teams and domains without rebuilding from scratch.",
        "Reduces discovery latency for code, notebooks, and platform operational know-how.",
        "Improves talent leverage by surfacing internal experts through vectorized profile intelligence.",
    ]),
    ("Platform Scope", [
        "Workspace intelligence: semantic retrieval over notebooks, scripts, and text assets.",
        "Knowledge intelligence: policy-grounded QA over platform documentation corpus.",
        "People intelligence: expertise discovery from generated profile embeddings.",
    ]),
    ("Architecture at a Glance", [
        "Experience layer: Next.js web portal with integrated assistant.",
        "Orchestration layer: FastAPI endpoints for search/profile/chat/admin sync.",
        "Intelligence layer: classification, rewrite, retrieval, and constrained generation.",
        "Data plane: ingestion catalog + Milvus collections + OpenAI model endpoints.",
    ]),
    ("Layered Design Principles", [
        "BFF abstraction decouples browser clients from backend schema volatility.",
        "Store-per-domain collections isolate retrieval workloads and simplify lifecycle management.",
        "Prompt externalization enables fast policy and behavior control without code redeploy.",
    ]),
    ("Ingestion Architecture", [
        "IngestionPipeline scans workspace directories and normalizes artifacts.",
        "Guardrails skip sensitive or unsupported files prior to indexing.",
        "Catalog + audit JSONs create deterministic source-of-truth for downstream jobs.",
    ]),
    ("Incremental Delta Strategy", [
        "Hash-based change tracking classifies NEW/UPDATED/UNCHANGED states.",
        "Incremental mode indexes only deltas, reducing infra spend and API token usage.",
        "Airflow DAG orchestrates recurring ingestion cadence for near-fresh knowledge.",
    ]),
    ("Vector Retrieval Stack", [
        "DocumentLoader assembles indexable content from ingestion catalog pointers.",
        "EmbeddingService generates vectors and caches payloads for efficiency.",
        "VectorRetriever/HybridRetriever serve semantic and hybrid search patterns.",
    ]),
    ("Generative Enrichment Design", [
        "Artifact summaries convert raw code context into explainable semantic capsules.",
        "Profile generation from summaries improves people-search signal quality.",
        "Summaries-first profile path minimizes prompt noise from raw artifacts.",
    ]),
    ("Chat Engine Architecture", [
        "Pipeline: intent classify -> query rewrite -> source retrieval -> response generation.",
        "Intent routing: DOC_QA, ARTIFACT_SEARCH, USER_SEARCH, HYBRID.",
        "USER_SEARCH favors deterministic name resolution before vector fallback.",
    ]),
    ("PromptOps Governance", [
        "All prompts live in versioned text assets under prompts/.",
        "Supports rapid policy tuning and controlled experimentation by prompt family.",
        "Separates orchestration code from conversational behavior artifacts.",
    ]),
    ("Frontend/BFF Architecture", [
        "Next.js App Router composes dashboard, workspaces, search, analytics, assistant.",
        "React Query standardizes fetch retries, staleness windows, and cache lifecycles.",
        "API routes normalize backend responses for stable UI contracts.",
    ]),
    ("Backend API Topology", [
        "Core endpoints: /query, /workspaces, /profile/workspace/{id}, /chat.",
        "Admin endpoints trigger indexing workflows for artifacts, profiles, docs.",
        "Health/metrics expose service state and operational telemetry.",
    ]),
    ("Data and Collection Model", [
        "kubeflow_artifacts: primary semantic retrieval corpus.",
        "artifact_summaries: compact generated abstraction layer for artifacts.",
        "user_profiles: expertise embeddings powering internal people search.",
        "platform_docs: documentation chunks for procedural QA.",
    ]),
    ("Runtime Flow: Search", [
        "Client search request -> Next API /api/search -> FastAPI /query.",
        "Retriever executes vector/hybrid query against kubeflow_artifacts collection.",
        "BFF remaps payload to UX-friendly card metadata schema.",
    ]),
    ("Runtime Flow: Chat", [
        "Client query/history -> /api/chat -> FastAPI /chat endpoint.",
        "ChatEngine executes intent detection and retrieval fan-out by route.",
        "LLM output returned with citations/artifacts/users for trust and explainability.",
    ]),
    ("Observability and Reliability", [
        "Startup preloads stores and degrades gracefully when optional stores are absent.",
        "Metrics include uptime, query volume, average latency, memory footprint, error rate.",
        "Health endpoint validates vector connectivity, model readiness, cache posture.",
    ]),
    ("Security and Guardrails", [
        "Content guardrails enforce skip logic for sensitive/unsupported artifacts.",
        "Server-side BFF proxy reduces direct backend exposure to browsers.",
        "Control points available for future policy engines and response moderation.",
    ]),
    ("Scalability Considerations", [
        "Independent collection scaling enables domain-specific performance tuning.",
        "Incremental indexing supports cost governance under growing data volume.",
        "Key watchpoints: embedding throughput, Milvus QPS/latency, prompt-token economics.",
    ]),
    ("GenAI Maturity Roadmap", [
        "Completed baseline: ingestion, retrieval API, web UX, summaries/profiles, chatbot.",
        "Next accelerators: LiteLLM gateway, Langfuse tracing, evaluation harnesses.",
        "Target state: production-grade LLMOps with SLOs and governed model lifecycle.",
    ]),
    ("Management Decision Themes", [
        "Platform is technically sound and extensible for enterprise GenAI expansion.",
        "Investment priority should shift to observability, controls, and deployment hardening.",
        "Recommend phase gate for production rollout with reliability and governance KPIs.",
    ]),
]


def add_brand_header_footer(slide) -> None:
    top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.35))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0.35), Inches(13.33), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35))
    tf = footer.text_frame
    tf.text = "Reliance GenAI Architecture Program"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def style_title(title_shape, color=BLUE, size=Pt(34)) -> None:
    title_shape.text_frame.paragraphs[0].font.color.rgb = color
    title_shape.text_frame.paragraphs[0].font.size = size
    title_shape.text_frame.paragraphs[0].font.bold = True


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_brand_header_footer(slide)

    slide.shapes.title.text = title
    style_title(slide.shapes.title, color=NAVY, size=Pt(32))

    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(21)
        p.font.color.rgb = DARK


def main() -> None:
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = WHITE
    add_brand_header_footer(title_slide)
    title_slide.shapes.title.text = TITLE
    style_title(title_slide.shapes.title, color=NAVY, size=Pt(40))
    title_slide.placeholders[1].text = SUBTITLE
    title_slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(22)
    title_slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = BLUE

    for title, bullets in slides_content:
        add_bullets_slide(prs, title, bullets)

    appendix = prs.slides.add_slide(prs.slide_layouts[5])
    appendix.background.fill.solid()
    appendix.background.fill.fore_color.rgb = LIGHT_BG
    add_brand_header_footer(appendix)
    appendix.shapes.title.text = "Appendix: Technical Module Map"
    style_title(appendix.shapes.title, color=NAVY, size=Pt(30))

    tx_box = appendix.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    lines = [
        "src/ingestion/: cli.py, pipeline.py, guards.py, extractors.py, storage.py",
        "src/retrieval/: api.py, indexer.py, retriever.py, embeddings.py, vector_store.py",
        "src/retrieval/chatbot/: engine.py, classifier.py, query_rewriter.py, retrievers.py, formatter.py",
        "webapp/app/api/: BFF proxy route handlers for all UI-to-backend interactions",
        "airflow/dags/ingestion_dag.py: scheduled incremental ingestion orchestration",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = DARK

    prs.save(OUTPUT_FILE)
    print(f"Created {OUTPUT_FILE} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
