"""
Single-slide capacity planning deck for the ingestion pipeline.
500 workspaces × 200 artifacts = 100K total artifacts.

Run:
    python generate_capacity_slide.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = "Pipeline_Capacity_Planning.pptx"

# ── Palette (matches v4) ──────────────────────────────────────────────────────
INK    = RGBColor(17,  24,  39)
MUTED  = RGBColor(75,  85,  99)
BLUE   = RGBColor(37,  99,  235)
GREEN  = RGBColor(22,  163, 74)
AMBER  = RGBColor(217, 119, 6)
PINK   = RGBColor(219, 39,  119)
TEAL   = RGBColor(13,  148, 136)
PURPLE = RGBColor(124, 58,  237)
RED    = RGBColor(220, 38,  38)
LIGHT  = RGBColor(248, 250, 252)
PANEL  = RGBColor(241, 245, 249)
WHITE  = RGBColor(255, 255, 255)
SLATE  = RGBColor(203, 213, 225)
WARN   = RGBColor(254, 226, 226)
GREENBG= RGBColor(220, 252, 231)
BLUEBG = RGBColor(219, 234, 254)
AMBERBG= RGBColor(254, 243, 199)
PINKBG = RGBColor(252, 231, 243)
PURPBG = RGBColor(237, 233, 254)


def add_textbox(slide, x, y, w, h, text, size=11, color=INK, bold=False,
                align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left  = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size   = Pt(size)
    p.font.bold   = bold
    p.font.italic = italic
    p.font.color.rgb = color
    return box


def set_text_lines(shape, lines, sizes, colors, bolds, aligns=None):
    """Multi-paragraph text in a shape."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left   = Inches(0.1)
    tf.margin_right  = Inches(0.08)
    tf.margin_top    = Inches(0.06)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (text, size, color, bold) in enumerate(zip(lines, sizes, colors, bolds)):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size  = Pt(size)
        p.font.color.rgb = color
        p.font.bold  = bold
        if aligns and aligns[i]:
            p.alignment = aligns[i]
        else:
            p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(1)


def add_box(slide, x, y, w, h, fill, line=SLATE, lw=0.75):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(lw)
    return shape


def add_rounded_box(slide, x, y, w, h, fill, line=SLATE, lw=0.75):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(lw)
    # corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def add_brand(slide):
    bar = add_box(slide, 0, 0, 13.33, 0.32, INK, INK)
    bar.line.fill.background()
    for i, col in enumerate([BLUE, GREEN, AMBER, PINK]):
        a = add_box(slide, i * 3.3325, 0.32, 3.3325, 0.07, col, col)
        a.line.fill.background()
    add_textbox(slide, 0.5, 7.1, 12.5, 0.28,
                "Kubeflow Workspace Intelligence LLMOps — Capacity Planning",
                size=8, color=MUTED, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# Main slide builder
# ─────────────────────────────────────────────────────────────────────────────
def make_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_brand(slide)

    # ── Title ────────────────────────────────────────────────────────────────
    add_textbox(slide, 0.45, 0.44, 8.8, 0.52,
                "Ingestion Pipeline — Capacity Planning at Scale",
                size=23, color=INK, bold=True)

    # ── Scale pills ──────────────────────────────────────────────────────────
    # "500 Workspaces" pill
    b1 = add_rounded_box(slide, 9.2, 0.44, 1.28, 0.46, BLUEBG, BLUE, 1.2)
    set_text_lines(b1, ["500", "Workspaces"],
                   [16, 8.5], [BLUE, INK], [True, False],
                   [PP_ALIGN.CENTER, PP_ALIGN.CENTER])
    add_textbox(slide, 10.5, 0.55, 0.22, 0.35, "×", 15, MUTED, True, PP_ALIGN.CENTER)
    b2 = add_rounded_box(slide, 10.74, 0.44, 1.28, 0.46, AMBERBG, AMBER, 1.2)
    set_text_lines(b2, ["~200", "Artifacts avg"],
                   [16, 8.5], [AMBER, INK], [True, False],
                   [PP_ALIGN.CENTER, PP_ALIGN.CENTER])
    add_textbox(slide, 12.04, 0.55, 0.22, 0.35, "=", 15, MUTED, True, PP_ALIGN.CENTER)
    b3 = add_rounded_box(slide, 12.28, 0.44, 0.85, 0.46, GREENBG, GREEN, 1.5)
    set_text_lines(b3, ["100K", "total"],
                   [14, 8.5], [GREEN, INK], [True, False],
                   [PP_ALIGN.CENTER, PP_ALIGN.CENTER])

    # ═════════════════════════════════════════════════════════════════════════
    # SECTION 1 — Pipeline Stage Timing
    # ═════════════════════════════════════════════════════════════════════════
    sec_label = add_box(slide, 0.45, 1.07, 12.45, 0.26, PANEL, SLATE, 0)
    sec_label.line.fill.background()
    add_textbox(slide, 0.55, 1.08, 4.0, 0.24,
                "STAGE-BY-STAGE TIME ESTIMATES  (Full Load: 100K artifacts)",
                size=8.5, color=MUTED, bold=True)

    # Stage definitions: (x, label, number, desc, seq, parallel, fill, line, seq_color, is_bottleneck)
    stages = [
        (0.45,  "01", "Data Ingestion",
         "Scan files, extract text,\nbuild catalog JSON",
         "~8 min", "~1–2 min",
         WHITE, SLATE, MUTED, False),
        (3.65,  "02", "Artifact Embedding",
         "OpenAI embed (batch 32)\n+ Milvus vector insert",
         "~1.8 hrs", "~8–10 min",
         BLUEBG, BLUE, BLUE, False),
        (6.85,  "03", "Artifact Summaries",
         "1 LLM call / artifact\ngpt-4o-mini, 220 tok max",
         "~75 hrs", "~3.5–4 hrs",
         WARN, RED, RED, True),
        (10.05, "04", "User Profiles",
         "1 LLM call / workspace\nfrom artifact summaries",
         "~35 min", "~5 min",
         PINKBG, PINK, PINK, False),
    ]

    box_w  = 2.92
    box_h  = 2.32
    top_y  = 1.40

    for (x, num, name, desc, seq, par, fill, lcolor, tcolor, is_bn) in stages:
        # outer box
        shape = add_rounded_box(slide, x, top_y, box_w, box_h, fill, lcolor, 1.2)

        # stage number badge
        badge = add_rounded_box(slide, x + 0.1, top_y + 0.09, 0.3, 0.27,
                                lcolor, lcolor, 0)
        badge.line.fill.background()
        add_textbox(slide, x + 0.1, top_y + 0.09, 0.3, 0.27,
                    num, 8.5, WHITE, True, PP_ALIGN.CENTER)

        # stage name
        add_textbox(slide, x + 0.46, top_y + 0.09, box_w - 0.56, 0.3,
                    name, 12, tcolor, True)

        # description
        add_textbox(slide, x + 0.12, top_y + 0.44, box_w - 0.22, 0.55,
                    desc, 9.5, MUTED, False)

        # divider line
        dl = add_box(slide, x + 0.12, top_y + 1.04, box_w - 0.24, 0.02,
                     SLATE, SLATE, 0)
        dl.line.fill.background()

        # sequential row
        add_textbox(slide, x + 0.12, top_y + 1.1, 1.3, 0.28,
                    "Sequential:", 8.5, MUTED, False)
        add_textbox(slide, x + 1.3, top_y + 1.1, box_w - 1.4, 0.28,
                    seq, 10.5, RED if is_bn else MUTED, True)

        # parallel row
        par_bg = add_rounded_box(slide, x + 0.1, top_y + 1.44, box_w - 0.2, 0.36,
                                 GREENBG, GREEN, 0.75)
        add_textbox(slide, x + 0.12, top_y + 1.48, 1.55, 0.28,
                    "20 workers:", 8.5, MUTED, False)
        add_textbox(slide, x + 1.55, top_y + 1.48, box_w - 1.65, 0.28,
                    par, 11, GREEN, True)

        # bottleneck badge
        if is_bn:
            bn = add_rounded_box(slide, x + 0.1, top_y + 1.88, box_w - 0.2, 0.30,
                                 WARN, RED, 0.75)
            add_textbox(slide, x + 0.12, top_y + 1.9, box_w - 0.22, 0.26,
                        "⚠  CRITICAL PATH — drives total runtime",
                        9, RED, True, PP_ALIGN.CENTER)

    # arrows between stages
    mid_y = top_y + box_h / 2
    for ax in [3.37, 6.57, 9.77]:
        add_arrow(slide, ax, mid_y, ax + 0.28, mid_y, BLUE, 2.0)
        # arrowhead triangle
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(ax + 0.20), Inches(mid_y - 0.075),
            Inches(0.14), Inches(0.15))
        tri.fill.solid()
        tri.fill.fore_color.rgb = BLUE
        tri.line.fill.background()
        tri.rotation = 90

    # ═════════════════════════════════════════════════════════════════════════
    # SECTION 2 — Execution Strategy
    # ═════════════════════════════════════════════════════════════════════════
    strat_y = top_y + box_h + 0.22

    sec2 = add_box(slide, 0.45, strat_y, 12.45, 0.26, PANEL, SLATE, 0)
    sec2.line.fill.background()
    add_textbox(slide, 0.55, strat_y + 0.01, 8.0, 0.24,
                "EXECUTION STRATEGY  — Full Load vs Incremental",
                size=8.5, color=MUTED, bold=True)

    bstrat_y = strat_y + 0.32

    # ── Full Load box ────────────────────────────────────────────────────────
    fl = add_rounded_box(slide, 0.45, bstrat_y, 5.88, 1.42, BLUEBG, BLUE, 1.5)

    add_textbox(slide, 0.58, bstrat_y + 0.07, 3.4, 0.28,
                "Full Load", 12.5, BLUE, True)

    rec1 = add_rounded_box(slide, 3.68, bstrat_y + 0.07, 2.52, 0.28,
                            BLUE, BLUE, 0)
    rec1.line.fill.background()
    add_textbox(slide, 3.68, bstrat_y + 0.07, 2.52, 0.28,
                "✓  BREADTH-FIRST (stage-by-stage)",
                8.5, WHITE, True, PP_ALIGN.CENTER)

    fl_points = [
        "Run all 500 workspaces through Stage 1, then Stage 2, then Stage 3, then Stage 4",
        "Maximises API batching & Airflow worker utilisation at each stage",
        "Stage-level checkpointing: resume from failed stage without restarting earlier stages",
    ]
    yy = bstrat_y + 0.40
    for pt in fl_points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.60), Inches(yy + 0.06), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        add_textbox(slide, 0.76, yy, 5.48, 0.29, pt, 9, INK)
        yy += 0.31

    # ── Incremental box ──────────────────────────────────────────────────────
    inc = add_rounded_box(slide, 6.54, bstrat_y, 5.88, 1.42, PINKBG, PINK, 1.5)

    add_textbox(slide, 6.68, bstrat_y + 0.07, 3.4, 0.28,
                "Incremental Load", 12.5, PINK, True)

    rec2 = add_rounded_box(slide, 9.78, bstrat_y + 0.07, 2.52, 0.28,
                            PINK, PINK, 0)
    rec2.line.fill.background()
    add_textbox(slide, 9.78, bstrat_y + 0.07, 2.52, 0.28,
                "✓  DEPTH-FIRST (per-workspace)",
                8.5, WHITE, True, PP_ALIGN.CENTER)

    inc_points = [
        "Each changed workspace flows through all 4 stages before the next workspace starts",
        "Per-workspace atomicity: consistent state, easier rollback on failure",
        "Content-hash guards skip unchanged artifacts — only deltas are re-embedded or re-summarised",
    ]
    yy = bstrat_y + 0.40
    for pt in inc_points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(6.70), Inches(yy + 0.06), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = PINK
        dot.line.fill.background()
        add_textbox(slide, 6.86, yy, 5.42, 0.29, pt, 9, INK)
        yy += 0.31

    # ── vs divider ───────────────────────────────────────────────────────────
    vs = add_rounded_box(slide, 6.16, bstrat_y + 0.50, 0.38, 0.36,
                          PANEL, SLATE, 0.75)
    add_textbox(slide, 6.16, bstrat_y + 0.50, 0.38, 0.36,
                "vs", 9, MUTED, True, PP_ALIGN.CENTER)

    # ═════════════════════════════════════════════════════════════════════════
    # SECTION 3 — Total Times + Assumptions
    # ═════════════════════════════════════════════════════════════════════════
    tot_y = bstrat_y + 1.50

    # Total times bar
    tb = add_rounded_box(slide, 0.45, tot_y, 12.45, 0.44, WHITE, SLATE, 0.75)

    add_textbox(slide, 0.60, tot_y + 0.06, 1.5, 0.30,
                "Total runtime:", 9.5, MUTED, True)

    t1 = add_rounded_box(slide, 2.18, tot_y + 0.06, 2.0, 0.30, WARN, RED, 0.75)
    add_textbox(slide, 2.18, tot_y + 0.06, 2.0, 0.30,
                "Sequential: ~75–80 hrs", 9, RED, True, PP_ALIGN.CENTER)

    add_textbox(slide, 4.3, tot_y + 0.06, 0.5, 0.30, "→", 11, MUTED, True, PP_ALIGN.CENTER)

    t2 = add_rounded_box(slide, 4.85, tot_y + 0.06, 2.55, 0.30, GREENBG, GREEN, 1.0)
    add_textbox(slide, 4.85, tot_y + 0.06, 2.55, 0.30,
                "20 workers:  ~4–4.5 hrs", 9.5, GREEN, True, PP_ALIGN.CENTER)

    add_textbox(slide, 7.55, tot_y + 0.06, 0.75, 0.30,
                "| Bottleneck:", 9, MUTED, False, PP_ALIGN.CENTER)

    t3 = add_rounded_box(slide, 8.4, tot_y + 0.06, 2.1, 0.30, WARN, RED, 0.75)
    add_textbox(slide, 8.4, tot_y + 0.06, 2.1, 0.30,
                "Stage 3 — Summaries", 9, RED, True, PP_ALIGN.CENTER)

    add_textbox(slide, 10.6, tot_y + 0.06, 2.2, 0.30,
                "↑ parallelise this stage first",
                9, MUTED, False, PP_ALIGN.RIGHT)

    # Assumptions
    add_textbox(slide, 0.45, tot_y + 0.52, 12.5, 0.24,
                "Assumptions:  OpenAI Tier 4 (gpt-4o-mini 10K RPM · text-embedding-3-small 5K RPM)  ·  "
                "20 Airflow workers  ·  batch_size = 32  ·  avg 200 artifacts/workspace  ·  "
                "LLM call latency ~2.5 s (Stage 3), ~2 s batch (Stage 2)",
                size=7.5, color=MUTED, italic=True)

    slide.notes_slide.notes_text_frame.text = (
        "Key numbers derived from observed run (2026-05-03): "
        "211 artifacts embedded in ~11s (6 batches), "
        "211 summaries in ~9.6 min (22/min sequential). "
        "Extrapolated to 100K artifacts. "
        "Stage 3 is the bottleneck: 100K LLM calls / 22 per min = 75 hrs sequential; "
        "with 20 parallel workers ~3.75 hrs. "
        "Stage 4 is fast because only 500 profile calls (one per workspace). "
        "Full load: breadth-first maximises batching and worker utilisation. "
        "Incremental: depth-first gives per-workspace atomicity and skips unchanged artifacts."
    )


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    make_slide(prs)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
