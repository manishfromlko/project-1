"""
5 slides — clean reference style (spacious, readable, dashed planes):
  1. Retrieval — Intent Classification
  2. Retrieval — Query Rewrite & Vector Search
  3. Generation — Context Assembly & LLM Answer
  4. Evaluation — Layer 1: Heuristic Scores
  5. Evaluation — Layer 2: RAGAS + Score Matrix

Run:  python generate_eval_slides_v2.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = "Retrieval_Generation_Evaluation_v2.pptx"

# ── Palette ────────────────────────────────────────────────────────────────────
WHITE  = RGBColor(255, 255, 255)
LIGHT  = RGBColor(250, 250, 252)
INK    = RGBColor(17,  24,  39)
MUTED  = RGBColor(100, 116, 139)
BLUE   = RGBColor(37,  99,  235)
GREEN  = RGBColor(22,  163, 74)
AMBER  = RGBColor(180, 100, 0)
RED    = RGBColor(185, 28,  28)
PINK   = RGBColor(190, 24,  93)
PURPLE = RGBColor(109, 40,  217)
TEAL   = RGBColor(15,  118, 110)
SLATE  = RGBColor(148, 163, 184)
DARKSL = RGBColor(71,  85,  105)
PANEL  = RGBColor(241, 245, 249)
BLUEBG = RGBColor(219, 234, 254)
GREENBG= RGBColor(220, 252, 231)
AMBERBG= RGBColor(254, 243, 199)
PINKBG = RGBColor(252, 231, 243)
PURPBG = RGBColor(237, 233, 254)
TEALB  = RGBColor(204, 251, 241)
WARN   = RGBColor(254, 226, 226)


# ── Primitives ─────────────────────────────────────────────────────────────────

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def rect(slide, x, y, w, h, fill=WHITE, line=SLATE, lw=1.0, dash=False, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(lw)
    if dash:
        s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if radius:
        s.adjustments[0] = 0.04
    return s


def txt(slide, x, y, w, h, text, size=10, color=INK, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    return b


def multi(slide, x, y, w, h, rows):
    """rows = [(text, size, color, bold), ...]"""
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (text, size, color, bold) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
    return b


def arrow(slide, x1, y1, x2, y2, color=DARKSL, lw=1.4):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    return c


def dot(slide, x, y, r=0.10, color=MUTED):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(r), Inches(r))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def bullets(slide, x, y, w, items, size=9, color=DARKSL, dot_color=MUTED, gap=0.28):
    """Render a bullet list; returns final y."""
    yy = y
    for item in items:
        dot(slide, x + 0.04, yy + 0.07, 0.09, dot_color)
        txt(slide, x + 0.20, yy, w - 0.24, gap, item, size, color)
        yy += gap
    return yy


def component_box(slide, x, y, w, h, title, fill, border, tsize=12, tcolor=None):
    """Solid colored component box with centered title — reference style."""
    rect(slide, x, y, w, h, fill, border, 1.5, radius=True)
    tc = tcolor or border
    txt(slide, x + 0.12, y + 0.10, w - 0.24, h - 0.20,
        title, tsize, tc, True, PP_ALIGN.CENTER)


def plane_box(slide, x, y, w, h, label, label_color=RED):
    """Dashed grouping boundary with plane label bottom-right."""
    rect(slide, x, y, w, h, fill=None, line=SLATE, lw=1.2, dash=True)
    txt(slide, x + w - 2.8, y + h - 0.30, 2.70, 0.28,
        label, 10, label_color, True, PP_ALIGN.RIGHT)


def slide_title(slide, text, sub=None):
    txt(slide, 0.30, 0.06, 12.8, 0.58, text, 22, INK, True)
    if sub:
        txt(slide, 0.32, 0.62, 12.7, 0.28, sub, 9.5, MUTED, False)


def bottom_bar(slide):
    for i, c in enumerate([BLUE, GREEN, AMBER, PINK]):
        rect(slide, i * 3.3325, 7.38, 3.3325, 0.12, c, c, 0)
    txt(slide, 0.5, 7.20, 12.5, 0.20,
        "Kubeflow Workspace Intelligence  ·  LLMOps",
        8, MUTED, False, PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Retrieval: Intent Classification
# ══════════════════════════════════════════════════════════════════════════════

def slide1_intent_classification(prs):
    slide = new_slide(prs)
    bottom_bar(slide)
    slide_title(slide,
        "Retrieval  –  Step 1: Intent Classification",
        "Every query is routed by an LLM before any vector search happens")

    # ── User query input ──────────────────────────────────────────────────────
    rect(slide, 3.90, 0.98, 5.50, 0.72, LIGHT, SLATE, 1.2)
    txt(slide, 4.0, 1.06, 5.3, 0.24, "User Query", 10, MUTED, True)
    txt(slide, 4.0, 1.28, 5.3, 0.32,
        '"How do I submit a Spark job?"', 10, INK, False, PP_ALIGN.CENTER)

    arrow(slide, 6.65, 1.70, 6.65, 1.98, DARKSL, 1.6)

    # ── Intent Classifier box ─────────────────────────────────────────────────
    rect(slide, 1.40, 1.98, 10.50, 1.76, BLUEBG, BLUE, 1.5)
    txt(slide, 1.60, 2.06, 5.0, 0.28, "Intent Classifier", 13, BLUE, True)
    txt(slide, 1.60, 2.34, 4.80, 0.22,
        "LLM call  ·  gpt-4o-mini  ·  temperature = 0.0", 9, DARKSL)

    # Detail bullets inside classifier
    bullets(slide, 1.60, 2.60, 4.40,
            ["System prompt defines 5 intent categories",
             "Returns JSON: { intent, confidence, reasoning }",
             "Falls back to DOC_QA on parse failure"],
            9, DARKSL, BLUE, 0.26)

    # JSON output example on right
    rect(slide, 7.40, 2.10, 4.10, 1.48, WHITE, BLUE, 0.8)
    multi(slide, 7.52, 2.16, 3.90, 1.36, [
        ("Returns:", 8.5, MUTED, True),
        ('  "intent":     "DOC_QA"', 9, INK, False),
        ('  "confidence": 0.92', 9, INK, False),
        ('  "reasoning":  "asks about platform', 9, INK, False),
        ('                 procedure"', 9, INK, False),
    ])

    arrow(slide, 6.65, 3.74, 6.65, 4.00, DARKSL, 1.6)

    # ── 5 intent outcome boxes ────────────────────────────────────────────────
    # 5 boxes × 2.28w + 4 × 0.14 gap = 11.4 + 0.56 = 11.96; left = 0.69
    intent_data = [
        (0.68,  "DOC_QA",           BLUEBG,  BLUE,
         ["Platform how-to guides", "Concepts & procedures", "Onboarding steps"]),
        (3.10,  "ARTIFACT_SEARCH",  GREENBG, GREEN,
         ["Find notebooks / scripts", "Code examples by tech", "Task-based search"]),
        (5.52,  "USER_SEARCH",      AMBERBG, AMBER,
         ["Find people by name", "Expertise / skill lookup", "Who works on what"]),
        (7.94,  "HYBRID",           PURPBG,  PURPLE,
         ["Query spans 2+ intents", "All 3 retrievers used", "top-3 from each store"]),
        (10.36, "OUT_OF_SCOPE",     WARN,    RED,
         ["Weather, stocks, news", "General coding help", "External services"]),
    ]

    for x, label, bg, fc, blist in intent_data:
        rect(slide, x, 4.00, 2.28, 2.70, bg, fc, 1.5, radius=True)
        txt(slide, x + 0.12, 4.08, 2.06, 0.30, label, 10, fc, True, PP_ALIGN.CENTER)
        rect(slide, x + 0.10, 4.40, 2.08, 0.02, fc, fc, 0.5)  # divider line
        bullets(slide, x + 0.10, 4.50, 2.10, blist, 8.5, DARKSL, fc, 0.26)

    # OUT_OF_SCOPE → immediate return note
    arrow(slide, 11.50, 6.70, 13.00, 6.70, RED, 1.4)
    rect(slide, 11.10, 6.72, 2.10, 0.54, WARN, RED, 0.8, radius=True)
    multi(slide, 11.18, 6.78, 1.94, 0.42, [
        ("⚡ Immediate return", 8.5, RED, True),
        ("No scoring · No retrieval", 8, DARKSL, False),
    ])

    # continuation note
    rect(slide, 0.68, 6.76, 9.60, 0.42, PANEL, SLATE, 0.6)
    txt(slide, 0.80, 6.82, 9.30, 0.30,
        "DOC_QA · ARTIFACT_SEARCH · USER_SEARCH · HYBRID  →  proceed to Step 2: Query Rewrite",
        9, DARKSL, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Retrieval: Query Rewrite & Vector Search
# ══════════════════════════════════════════════════════════════════════════════

def slide2_vector_search(prs):
    slide = new_slide(prs)
    bottom_bar(slide)
    slide_title(slide,
        "Retrieval  –  Step 2: Query Rewrite & Vector Search",
        "Query is enriched before embedding, then routed to the correct Milvus collection")

    # ── Query Rewriter ────────────────────────────────────────────────────────
    rect(slide, 2.40, 0.98, 8.50, 1.54, TEALB, TEAL, 1.5)
    txt(slide, 2.60, 1.06, 4.0, 0.28, "Query Rewriter", 13, TEAL, True)
    txt(slide, 2.60, 1.34, 3.90, 0.22,
        "LLM call  ·  gpt-4o-mini  ·  T=0.0  ·  max 80 tokens", 9, DARKSL)
    bullets(slide, 2.60, 1.58, 3.80,
            ["Expands query with synonyms and technical terms",
             "Improves cosine similarity recall in Milvus"],
            9, DARKSL, TEAL, 0.26)

    # example rewrite on right of rewriter box
    rect(slide, 6.80, 1.10, 3.80, 1.28, WHITE, TEAL, 0.8)
    multi(slide, 6.92, 1.16, 3.60, 1.18, [
        ("Example:", 8.5, MUTED, True),
        ('In:   "Spark job Kubeflow"', 9, DARKSL, False),
        ('Out:  "submit Spark application Kubeflow', 9, TEAL, False),
        ('       pipeline resource configuration"', 9, TEAL, False),
    ])

    # arrows down from rewriter to 3 retriever columns
    arrow(slide, 3.50, 2.52, 3.50, 2.86, GREEN,  1.4)
    arrow(slide, 6.65, 2.52, 6.65, 2.86, AMBER,  1.4)
    arrow(slide, 9.80, 2.52, 9.80, 2.86, PURPLE, 1.4)

    # ── Three retriever columns ───────────────────────────────────────────────
    col_data = [
        (1.62, "Doc Retriever",       GREEN,  GREENBG,
         "platform_docs",
         ["Kubeflow access guide", "Notebook setup guide", "Spark job guide"],
         ["chunk_text  (top 5)", "HNSW · COSINE · 1536-dim", "Used by:  DOC_QA · HYBRID"]),
        (4.77, "Artifact Retriever",  AMBER,  AMBERBG,
         "artifact_summaries",
         ["ajay11.yadav — PySpark", "amit23.sharma — pandas", "dhruv2.aggarwal — Spark"],
         ["artifact_summary  (top 5)", "HNSW · COSINE · 1536-dim", "Used by:  ARTIFACT · HYBRID"]),
        (7.92, "User Retriever",      PURPLE, PURPBG,
         "user_profiles",
         ["ravi.verma", "julie1.dixit", "+ 3 more profiles"],
         ["user_profile  (top 5)", "HNSW · COSINE · 1536-dim", "Used by:  USER · HYBRID"]),
    ]

    for x, title, fc, bg, coll, examples, detail in col_data:
        # retriever box
        rect(slide, x, 2.86, 3.60, 0.70, bg, fc, 1.5, radius=True)
        txt(slide, x + 0.14, 2.96, 3.34, 0.26, title, 11.5, fc, True, PP_ALIGN.CENTER)
        txt(slide, x + 0.14, 3.24, 3.34, 0.22, coll, 9, DARKSL, False, PP_ALIGN.CENTER)

        arrow(slide, x + 1.80, 3.56, x + 1.80, 3.82, fc, 1.4)

        # Milvus collection box
        rect(slide, x, 3.82, 3.60, 2.12, WHITE, fc, 1.2)
        txt(slide, x + 0.14, 3.90, 3.34, 0.26, "Milvus Collection", 9, MUTED, True)
        yy = 4.20
        for ex in examples:
            dot(slide, x + 0.18, yy + 0.06, 0.09, fc)
            txt(slide, x + 0.34, yy, 3.10, 0.26, ex, 9, DARKSL)
            yy += 0.28

        # detail strip at bottom of collection box
        rect(slide, x, 5.94, 3.60, 0.58, PANEL, fc, 0.6)
        yy = 6.00
        for d in detail:
            txt(slide, x + 0.14, yy, 3.34, 0.19, d, 8, MUTED, False)
            yy += 0.19

    # ── Name Resolution path (USER_SEARCH special) ────────────────────────────
    plane_box(slide, 0.20, 6.58, 13.0, 0.68, "Special path:  USER_SEARCH exact name", PURPLE)
    multi(slide, 0.36, 6.64, 12.60, 0.56, [
        ("Name Resolution:  ", 9, PURPLE, True),
        ("RapidFuzz string match against all user IDs  →  LLM disambiguation (user_resolve span)  →  "
         "direct Milvus profile fetch  →  LLM judge scores profile_relevance.  "
         "No generate step, no Layer 1 scores.", 9, DARKSL, False),
    ])
    # make it inline
    b = slide.shapes.add_textbox(Inches(0.36), Inches(6.64), Inches(12.60), Inches(0.52))
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.font.size = Pt(9); p.font.color.rgb = DARKSL; p.alignment = PP_ALIGN.LEFT
    run1 = p.add_run(); run1.text = "Name Resolution:  "
    run1.font.bold = True; run1.font.color.rgb = PURPLE; run1.font.size = Pt(9)
    run2 = p.add_run()
    run2.text = ("RapidFuzz string match against all user IDs  →  LLM disambiguation (user_resolve span)  →  "
                 "direct Milvus profile fetch  →  LLM judge scores profile_relevance.  "
                 "No generate step. No Layer 1 heuristic scores.")
    run2.font.color.rgb = DARKSL; run2.font.size = Pt(9)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 — Generation: Context Assembly & LLM Answer
# ══════════════════════════════════════════════════════════════════════════════

def slide3_generation(prs):
    slide = new_slide(prs)
    bottom_bar(slide)
    slide_title(slide,
        "Generation  –  Context Assembly & LLM Answer",
        "Retrieved chunks are assembled into a prompt and sent to gpt-4o-mini via LiteLLM proxy")

    # ── Step boxes (left to right) ────────────────────────────────────────────
    steps = [
        (0.28, "Retrieved\nContext", GREENBG, GREEN),
        (3.66, "Prompt\nBuilder",    BLUEBG,  BLUE),
        (7.04, "LLM\nGenerate",      PURPBG,  PURPLE),
        (10.42,"Response\nSchema",   AMBERBG, AMBER),
    ]
    for x, label, bg, fc in steps:
        rect(slide, x, 0.98, 3.10, 0.80, bg, fc, 1.5, radius=True)
        txt(slide, x + 0.14, 1.05, 2.84, 0.66,
            label, 13, fc, True, PP_ALIGN.CENTER)

    for x in [3.38, 6.76, 10.14]:
        arrow(slide, x, 1.38, x + 0.28, 1.38, DARKSL, 1.6)

    # ── Detail sections below each step ──────────────────────────────────────

    # Retrieved Context details
    rect(slide, 0.28, 1.90, 3.10, 2.20, WHITE, GREEN, 1.0)
    txt(slide, 0.42, 1.98, 2.84, 0.24, "Per Intent", 9, MUTED, True)
    detail_ctx = [
        ("DOC_QA",          "doc chunks   (top 5)"),
        ("ARTIFACT_SEARCH", "summaries     (top 5)"),
        ("USER_SEARCH",     "profiles       (top 5)"),
        ("HYBRID",          "all 3 sources (top 3 each)"),
    ]
    yy = 2.26
    for intent, src in detail_ctx:
        rect(slide, 0.38, yy, 3.00, 0.40, PANEL, SLATE, 0.5)
        txt(slide, 0.50, yy + 0.04, 1.10, 0.32, intent, 8.5, DARKSL, True)
        txt(slide, 1.62, yy + 0.04, 1.68, 0.32, src, 8.5, MUTED, False)
        yy += 0.44

    # Prompt Builder details
    rect(slide, 3.66, 1.90, 3.10, 2.20, WHITE, BLUE, 1.0)
    txt(slide, 3.80, 1.98, 2.84, 0.24, "4 Templates", 9, MUTED, True)
    prompts = [
        ("doc_qa",          "answer only from retrieved docs"),
        ("artifact_search", "list matching artifact IDs"),
        ("user_search",     "describe most relevant person"),
        ("hybrid",          "combine all 3 context types"),
    ]
    yy = 2.26
    for tmpl, desc in prompts:
        rect(slide, 3.76, yy, 2.90, 0.40, PANEL, SLATE, 0.5)
        txt(slide, 3.88, yy + 0.04, 1.00, 0.32, tmpl, 8.5, BLUE, True)
        txt(slide, 4.90, yy + 0.04, 1.68, 0.32, desc, 8, MUTED, False)
        yy += 0.44
    txt(slide, 3.80, 4.16, 2.90, 0.22,
        "Chat history prepended between system + user msg", 8, MUTED, False)

    # LLM Generate details
    rect(slide, 7.04, 1.90, 3.10, 2.20, WHITE, PURPLE, 1.0)
    txt(slide, 7.18, 1.98, 2.84, 0.24, "Parameters", 9, MUTED, True)
    llm_params = [
        ("Model",       "gpt-4o-mini"),
        ("Temperature", "0.2"),
        ("Max tokens",  "600"),
        ("Routing",     "via LiteLLM proxy"),
        ("Trace name",  "chat · {intent}"),
        ("Tags",        "intent:X · model:X"),
    ]
    yy = 2.26
    for k, v in llm_params:
        txt(slide, 7.18, yy, 1.14, 0.26, k, 8.5, DARKSL, True)
        txt(slide, 8.34, yy, 1.72, 0.26, v, 8.5, MUTED, False)
        yy += 0.30

    # Response schema
    rect(slide, 10.42, 1.90, 2.68, 2.20, WHITE, AMBER, 1.0)
    txt(slide, 10.56, 1.98, 2.40, 0.24, "Response fields", 9, MUTED, True)
    fields = ["answer  (str)",
              "intent  (str)",
              "confidence  (float)",
              "trace_id  (uuid)",
              "raw_docs  (list)",
              "raw_artifacts  (list)",
              "raw_users  (list)"]
    bullets(slide, 10.52, 2.26, 2.56, fields, 8.5, DARKSL, AMBER, 0.26)

    # ── Langfuse Trace section ────────────────────────────────────────────────
    plane_box(slide, 0.28, 4.22, 12.80, 2.60, "Langfuse Trace (per request)", BLUE)

    txt(slide, 0.44, 4.30, 5.50, 0.26,
        "trace_id  —  UUID generated at request start, forwarded via LiteLLM extra_body to group all spans",
        9, DARKSL, False)

    # Trace spans
    span_data = [
        (0.44,  "classify",       PURPBG,  PURPLE, "gpt-4o-mini · T=0.0 · max 150 tok\nOutputs: intent, confidence, reasoning"),
        (3.56,  "rewrite",        TEALB,   TEAL,   "gpt-4o-mini · T=0.0 · max 80 tok\nOutputs: enriched search query string"),
        (6.68,  "generate",       BLUEBG,  BLUE,   "gpt-4o-mini · T=0.2 · max 600 tok\nMetadata: query · intent · conf · hits"),
        (9.80,  "user_resolve",   AMBERBG, AMBER,  "gpt-4o-mini · T=0.0 (USER_SEARCH only)\nDisambiguates fuzzy name matches"),
    ]
    for x, name, bg, fc, desc in span_data:
        rect(slide, x, 4.64, 2.96, 1.38, bg, fc, 1.2, radius=True)
        txt(slide, x + 0.14, 4.72, 2.70, 0.28, name, 10.5, fc, True, PP_ALIGN.CENTER)
        rect(slide, x + 0.12, 5.02, 2.72, 0.02, fc, fc, 0.5)
        txt(slide, x + 0.14, 5.08, 2.70, 0.84, desc, 8.5, DARKSL, False)

    txt(slide, 0.44, 6.08, 12.60, 0.28,
        "Scores are posted to Langfuse AFTER the trace via score_response_quality() and evaluate_in_background() — "
        "not attached to individual spans but to the root trace by trace_id.",
        8.5, MUTED, False)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 — Evaluation: Layer 1 Heuristic Scores
# ══════════════════════════════════════════════════════════════════════════════

def slide4_layer1(prs):
    slide = new_slide(prs)
    bottom_bar(slide)
    slide_title(slide,
        "Evaluation  –  Layer 1: Heuristic Scores",
        "4 scores posted synchronously inline after every generate call  ·  no LLM involved  ·  appear in Langfuse within 1–2 seconds")

    # ── 4 large score cards (2 × 2 grid) ─────────────────────────────────────
    cards = [
        # (x, y, name, color, bg, formula_title, formula_body, good, bad)
        (0.28, 0.95,
         "response_length", BLUE, BLUEBG,
         "Normalised answer length",
         "< 50 chars     → ramps from 0.0 to 0.5\n"
         "50 – 300 chars  → ramps from 0.5 to 1.0\n"
         "300 – 2000 chars → 1.0  (sweet spot)\n"
         "> 2000 chars    → penalised toward 0.5",
         "Answer is well-developed and substantive",
         "Very short fallback or extremely verbose wall of text"),

        (6.90, 0.95,
         "has_content", GREEN, GREENBG,
         "Substantive answer check  (string match, no LLM)",
         'Scans final answer for fallback phrases:\n'
         '"i couldn\'t find"  ·  "no matching"  ·  "please try again"\n'
         '"i don\'t have access"  ·  "i encountered an error"\n\n'
         "1.0  if none found   ·   0.0  if any found",
         "LLM answered from retrieved context",
         "LLM gave up — retrieved nothing relevant"),

        (0.28, 4.00,
         "intent_confidence", PURPLE, PURPBG,
         "Classifier certainty  (passthrough, no transform)",
         "Direct passthrough of confidence value\nreturned by the Intent Classifier LLM call.\n\n"
         "High (> 0.80)  →  clear, unambiguous query\n"
         "Medium (0.60–0.80)  →  boundary query\n"
         "Low (< 0.60)  →  classifier was uncertain",
         "Query matched a clear intent category",
         "Ambiguous query or classifier prompt needs tuning"),

        (6.90, 4.00,
         "source_count", AMBER, AMBERBG,
         "Retrieval coverage",
         "value  =  (doc_hits + artifact_hits + user_hits)  ÷  5\n"
         "Capped at 1.0\n\n"
         "5 sources → 1.0   |   3 sources → 0.6\n"
         "1 source  → 0.2   |   0 sources → 0.0",
         "Multiple relevant sources found in Milvus",
         "Retrieval returned nothing — embedding or indexing issue"),
    ]

    for x, y, name, fc, bg, formula_title, formula_body, good, bad in cards:
        # Card border
        rect(slide, x, y, 6.28, 2.84, bg, fc, 1.8, radius=True)

        # Score name header
        rect(slide, x + 0.14, y + 0.14, 6.00, 0.46, fc, fc, 0, radius=True)
        txt(slide, x + 0.24, y + 0.18, 5.80, 0.38, name, 14, WHITE, True, PP_ALIGN.CENTER)

        # Formula title
        txt(slide, x + 0.20, y + 0.70, 5.90, 0.26, formula_title, 9.5, INK, True)

        # Formula body
        txt(slide, x + 0.20, y + 0.96, 5.90, 0.92, formula_body, 9, DARKSL, False)

        # Good / Bad signals
        rect(slide, x + 0.14, y + 2.00, 2.90, 0.68, WHITE, GREEN, 0.6)
        txt(slide, x + 0.24, y + 2.06, 2.70, 0.18, "Green signal", 7.5, GREEN, True)
        txt(slide, x + 0.24, y + 2.24, 2.70, 0.36, good, 8, DARKSL, False)

        rect(slide, x + 3.22, y + 2.00, 2.92, 0.68, WHITE, RED, 0.6)
        txt(slide, x + 3.32, y + 2.06, 2.72, 0.18, "Red signal", 7.5, RED, True)
        txt(slide, x + 3.32, y + 2.24, 2.72, 0.36, bad, 8, DARKSL, False)

    # ── NOT posted warning ────────────────────────────────────────────────────
    rect(slide, 0.28, 6.92, 12.78, 0.36, WARN, RED, 1.0)
    b = slide.shapes.add_textbox(Inches(0.44), Inches(6.96), Inches(12.44), Inches(0.28))
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = "⚠  Not posted for:  "
    r1.font.bold = True; r1.font.size = Pt(9); r1.font.color.rgb = RED
    r2 = p.add_run()
    r2.text = ("OUT_OF_SCOPE  (engine short-circuits before score block)   ·   "
               "Exact USER_SEARCH name match  (engine returns before score block via early-return on exact_uid)")
    r2.font.size = Pt(9); r2.font.color.rgb = DARKSL


# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 — Evaluation: Layer 2 RAGAS + Score Matrix
# ══════════════════════════════════════════════════════════════════════════════

def slide5_layer2_matrix(prs):
    slide = new_slide(prs)
    bottom_bar(slide)
    slide_title(slide,
        "Evaluation  –  Layer 2: RAGAS Scores & Score Matrix",
        "Run in a background daemon thread  ·  results appear in Langfuse 10–30 seconds after the trace")

    # ── Layer 2 score cards (top row) ─────────────────────────────────────────
    l2_cards = [
        (0.28,  "faithfulness",      PURPLE, PURPBG,
         "Answer grounded in context",
         "RAGAS metric  ·  eval model: gpt-4o\n"
         "Measures: is every claim in the\n"
         "answer supported by the retrieved\n"
         "chunks / summaries?\n\n"
         "Low → LLM hallucinated from training"),

        (3.56,  "answer_relevancy",  BLUE,   BLUEBG,
         "Answer addresses the question",
         "RAGAS metric  ·  uses embeddings\n"
         "Measures: does the answer actually\n"
         "respond to what was asked?\n\n"
         "Low → answer is off-topic or vague"),

        (6.84,  "context_precision", TEAL,   TEALB,
         "Top context is most relevant",
         "RAGAS metric  ·  eval model: gpt-4o\n"
         "Measures: are the highest-ranked\n"
         "chunks more relevant than lower ones?\n\n"
         "Low → wrong chunks ranked first"),

        (10.12, "profile_relevance", PINK,   PINKBG,
         "Profile answers the question",
         "LLM judge  ·  eval model: gpt-4o\n"
         "Scores: 1.0 / 0.7 / 0.3 / 0.0\n"
         "Only for exact USER_SEARCH match\n"
         "(profile fetched directly, no generate)\n\n"
         "Low → wrong user resolved"),
    ]

    for x, name, fc, bg, title, desc in l2_cards:
        rect(slide, x, 0.95, 3.04, 3.08, bg, fc, 1.5, radius=True)
        rect(slide, x + 0.12, 1.05, 2.80, 0.42, fc, fc, 0, radius=True)
        txt(slide, x + 0.18, 1.09, 2.68, 0.34, name, 11, WHITE, True, PP_ALIGN.CENTER)
        txt(slide, x + 0.16, 1.56, 2.76, 0.24, title, 9, INK, True)
        txt(slide, x + 0.16, 1.80, 2.76, 1.14, desc, 9, DARKSL, False)

    # user feedback smaller card (appended right of profile_relevance)
    # Actually fits better as a note below

    # ── User feedback note ─────────────────────────────────────────────────────
    rect(slide, 0.28, 4.14, 3.04, 1.20, PINKBG, PINK, 1.2, radius=True)
    txt(slide, 0.44, 4.22, 2.72, 0.28, "user_feedback", 11, PINK, True)
    txt(slide, 0.44, 4.50, 2.72, 0.22, "User-initiated via frontend thumbs", 9, DARKSL, False)
    multi(slide, 0.44, 4.72, 2.72, 0.52, [
        ("👍  Thumbs-up   →  1.0", 9, DARKSL, False),
        ("👎  Thumbs-down →  0.0", 9, DARKSL, False),
        ("POST /observability/feedback", 8.5, MUTED, False),
    ])

    # ── Score matrix ───────────────────────────────────────────────────────────
    txt(slide, 3.56, 4.14, 9.76, 0.30,
        "Score Availability Matrix", 11, INK, True)

    # Column headers
    col_xs =  [3.56, 5.60, 7.32, 8.92, 10.52, 12.14]
    col_ws =  [1.96, 1.64, 1.52, 1.52, 1.54,  1.06]
    col_hdr = ["Intent Path",
               "Layer 1\n(4 scores)",
               "faithfulness\nanswer_rel.\nctx_precision",
               "profile\nrelevance",
               "user\nfeedback",
               "LLM\ncalls"]

    for cx, cw, hdr in zip(col_xs, col_ws, col_hdr):
        rect(slide, cx, 4.48, cw - 0.04, 0.66, INK, INK, 0)
        txt(slide, cx + 0.06, 4.52, cw - 0.14, 0.60, hdr, 8, WHITE, True, PP_ALIGN.CENTER)

    # Matrix rows
    Y = "✓"
    N = "–"
    rows = [
        ("DOC_QA",               BLUEBG,  BLUE,   Y, Y, N, Y, "3"),
        ("ARTIFACT_SEARCH",      GREENBG, GREEN,  Y, Y, N, Y, "3"),
        ("USER_SEARCH (semantic)",AMBERBG, AMBER,  Y, Y, N, Y, "3"),
        ("USER_SEARCH (exact)",  PINKBG,  PINK,   N, N, Y, Y, "3"),
        ("HYBRID",               PURPBG,  PURPLE, Y, Y, N, Y, "3"),
        ("OUT_OF_SCOPE",         WARN,    RED,    N, N, N, N, "1"),
    ]
    for ri, (name, bg, fc, l1, ragas, prof, fb, calls) in enumerate(rows):
        ry = 5.18 + ri * 0.30
        cells = [name, l1, ragas, prof, fb, calls]
        for ci, (cx, cw, cell) in enumerate(zip(col_xs, col_ws, cells)):
            cell_bg = bg if ci == 0 else WHITE
            rect(slide, cx, ry, cw - 0.04, 0.28, cell_bg, SLATE, 0.4)
            fc_val = (GREEN if cell == Y else (RED if cell == N else INK)) if ci > 0 else fc
            txt(slide, cx + 0.06, ry + 0.03, cw - 0.14, 0.24,
                cell, 9 if ci == 0 else 10, fc_val,
                ci == 0, PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)

    # Threshold legend
    rect(slide, 3.56, 6.98, 9.76, 0.30, PANEL, SLATE, 0.5)
    b = slide.shapes.add_textbox(Inches(3.64), Inches(7.02), Inches(9.60), Inches(0.24))
    tf = b.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.04); tf.margin_top = Inches(0.01)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    thresholds = [
        ("Green: ", GREEN),  ("faithfulness > 0.70  |  answer_relevancy > 0.75  |  intent_confidence > 0.80  |  has_content = 1.0   ", DARKSL),
        ("Red: ", RED),      ("faithfulness < 0.50  |  has_content = 0.0  |  source_count = 0.0", DARKSL),
    ]
    for text, color in thresholds:
        run = p.add_run(); run.text = text
        run.font.size = Pt(8); run.font.color.rgb = color
        run.font.bold = (color in (GREEN, RED))


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide1_intent_classification(prs)
    slide2_vector_search(prs)
    slide3_generation(prs)
    slide4_layer1(prs)
    slide5_layer2_matrix(prs)

    prs.save(OUTPUT)
    print(f"Saved → {OUTPUT}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
